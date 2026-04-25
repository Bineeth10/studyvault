# =====================================================
# SECTION: Imports
# Purpose: Import all required libraries and modules
# =====================================================
from fastapi import APIRouter, Request, Form, UploadFile, File
from fastapi.responses import HTMLResponse, RedirectResponse, JSONResponse, FileResponse
from fastapi.templating import Jinja2Templates
from datetime import datetime, timezone
from bson import ObjectId
import os
import shutil
from typing import List, Optional, Any
import asyncio
import json
import re
import concurrent.futures
import difflib
from groq import Groq
from dotenv import load_dotenv
import PyPDF2
from PIL import Image
import pytesseract
import docx
from pptx import Presentation

from app.database import (
    users_collection, notes_collection, favorites_collection,
    notifications_collection, messages_collection, db,
    create_notification, system_settings_collection, reports_collection,
    announcements_collection, subjects_collection, forum_posts_collection
)
from app.templates import templates

# =====================================================
# SECTION: Configuration
# Purpose: Initialize the router and storage settings
# =====================================================
router = APIRouter(prefix="/student")

# Local storage for uploaded files
UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# =====================================================
# SECTION: Utility Functions
# Purpose: Helper functions for text extraction and context
# =====================================================

def extract_text_from_pdf(file_path: str) -> str:
    # Extracts plain text from PDF files using PyPDF2
    # Returns the first 3000 characters for processing
    pages_text: list = []
    try:
        with open(file_path, 'rb') as pdf_file:
            try:
                pdf_reader = PyPDF2.PdfReader(pdf_file)
                running_len: int = 0
                for page in pdf_reader.pages:
                    raw = page.extract_text()
                    page_text: str = raw if isinstance(raw, str) else ""
                    if page_text:
                        pages_text.append(page_text)
                        running_len = running_len + len(page_text)  # type: ignore[operator]
                    if running_len >= 3000:
                        break
            except Exception as e:
                print(f"PyPDF2 reading error: {e}")
    except Exception as e:
        print(f"Error opening PDF: {e}")

    extracted_text: str = "\n".join(pages_text)
    if not extracted_text or len(extracted_text.strip()) < 50:
        return ""

    return extracted_text[:3000]  # type: ignore[index]

def extract_text_from_image(file_path: str) -> str:
    # Uses Tesseract OCR to extract text from image files (JPG, PNG, etc.)
    try:
        if os.name == 'nt':
            tesseract_path = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
            if os.path.exists(tesseract_path):
                pytesseract.pytesseract.tesseract_cmd = tesseract_path

        img = Image.open(file_path)
        raw_text = pytesseract.image_to_string(img)
        text: str = raw_text if isinstance(raw_text, str) else ""

        if not text or len(text.strip()) < 10:
            return ""

        return text[:4000]  # type: ignore[index]
    except Exception as e:
        print(f"⚠️ Error extracting text from image: {e}")
        return ""

def extract_text_from_office(file_path: str) -> str:
    # Extracts text from modern MS Office files (DOCX, PPTX)
    ext = os.path.splitext(file_path)[1].lower()
    text: str = ""
    try:
        if ext == ".docx":
            doc = docx.Document(file_path)
            text = "\n".join([str(p.text) for p in doc.paragraphs])
        elif ext == ".pptx":
            prs = Presentation(file_path)
            parts: list = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text: str = str(getattr(shape, "text", ""))
                        if shape_text:
                            parts.append(shape_text)
            text = "\n".join(parts)

        if not text or len(text.strip()) < 10:
            return ""

        return text[:4000]  # type: ignore[index]
    except Exception as e:
        print(f"Office extraction error ({ext}): {e}")
        return ""

def get_notif_context(user_id: str) -> dict:
    # Fetches unread notification count and recent alerts for the layout
    try:
        unread_count = notifications_collection.count_documents({"user_id": user_id, "is_read": False})
        recent_notifications = list(
            notifications_collection.find({"user_id": user_id}).sort("created_at", -1).limit(5)
        )
    except Exception:
        unread_count = 0
        recent_notifications = []
    return {"unread_count": unread_count, "recent_notifications": recent_notifications}


def check_duplicate_and_plagiarism(
    title: str, extracted_text: str, uploader_id: str,
    subject: str = "", similarity_threshold: float = 0.75
) -> dict:
    """
    Checks the uploaded note against existing approved/pending notes in the DB.
    Returns a dict with keys: status, plagiarism_score, duplicate_of, detail
    
    Status values:
      'Safe'                - unique title + content
      'Duplicate Title'     - exact title match found
      'Possible Plagiarism' - content similarity >= threshold
    """
    title_norm = title.strip().lower()

    # ── Step 1: Exact title match (any subject, any uploader except self) ──────
    existing_same_title = notes_collection.find_one({
        "title": {"$regex": f"^{re.escape(title)}$", "$options": "i"},
        "uploader_id": {"$ne": str(uploader_id)}
    })
    if existing_same_title:
        print(f"[Plagiarism] Duplicate title detected: '{title}' matches note {existing_same_title['_id']}")
        return {
            "status": "Duplicate Title",
            "plagiarism_score": 1.0,
            "duplicate_of": str(existing_same_title["_id"]),
            "duplicate_uploader": existing_same_title.get("uploader_name", "Unknown"),
            "detail": f"Title '{title}' was already uploaded by {existing_same_title.get('uploader_name', 'another student')}."
        }

    # ── Step 2: Text similarity (only if we have extracted text) ──────────────
    if not extracted_text or len(extracted_text.strip()) < 100:
        # Not enough text to meaningfully compare
        return {"status": "Safe", "plagiarism_score": 0.0, "duplicate_of": None, "detail": ""}

    # Fetch recent notes in the same subject (limit to 30 for performance)
    query: dict = {"uploader_id": {"$ne": str(uploader_id)}}
    if subject:
        query["subject"] = subject
    candidates = list(notes_collection.find(query).sort("uploaded_at", -1).limit(30))

    best_score = 0.0
    best_match = None
    new_text_lower = extracted_text.lower()

    for candidate in candidates:
        cand_path = candidate.get("file_path", "")
        if not cand_path or not os.path.exists(cand_path):
            continue
        try:
            # Extract text from candidate
            fn = (candidate.get("filename") or "").lower()
            cand_text = ""
            if fn.endswith(".pdf"):
                cand_text = extract_text_from_pdf(cand_path)
            elif any(fn.endswith(ext) for ext in [".docx", ".pptx"]):
                cand_text = extract_text_from_office(cand_path)

            if not cand_text or len(cand_text.strip()) < 100:
                continue

            # SequenceMatcher ratio: 0.0 (no match) → 1.0 (identical)
            ratio = difflib.SequenceMatcher(
                None,
                new_text_lower[:3000],
                cand_text.lower()[:3000],
                autojunk=True
            ).ratio()

            print(f"[Plagiarism] Similarity vs '{candidate.get('title')}': {ratio:.2%}")

            if ratio > best_score:
                best_score = ratio
                best_match = candidate
        except Exception as e:
            print(f"[Plagiarism] Error comparing with candidate {candidate.get('_id')}: {e}")
            continue

    if best_score >= similarity_threshold and best_match:
        print(f"[Plagiarism] Flagged: {best_score:.2%} similarity with '{best_match.get('title')}'")
        return {
            "status": "Possible Plagiarism",
            "plagiarism_score": round(best_score, 4),
            "duplicate_of": str(best_match["_id"]),
            "duplicate_uploader": best_match.get("uploader_name", "Unknown"),
            "detail": f"{best_score:.0%} content similarity with note '{best_match.get('title')}' by {best_match.get('uploader_name', 'another student')}."
        }

    print(f"[Plagiarism] Safe. Best similarity was {best_score:.2%}")
    return {"status": "Safe", "plagiarism_score": round(best_score, 4), "duplicate_of": None, "detail": ""}

# =====================================================
# SECTION: Authentication Logic
# Purpose: Handles role-based session validation
# =====================================================

def get_current_student(request: Request):
    # Dependency to ensure the user is logged in as a student
    user = request.session.get("user")
    if not user or user.get("role") != "student":
        return None
    
    # Refresh user data from DB to ensure session data is current
    full_user = users_collection.find_one({"_id": ObjectId(user["id"])})
    if full_user:
        user = user.copy()
        user.update(full_user)
        user["id"] = str(full_user["_id"])
        
    return user

def safe_id(id_val):
    """Safely converts string to ObjectId, or returns original if invalid hex."""
    if not id_val: return None
    if isinstance(id_val, ObjectId): return id_val
    try: return ObjectId(str(id_val))
    except: return str(id_val)

# =====================================================
# SECTION: AI Processing
# Purpose: Smart summary, safety analysis, and demo modes
# =====================================================

def generate_ai_summary(text):
    # Generates a 3-6 sentence academic summary using Groq Llama-3.1
    if not text or not text.strip():
        return None
    try:
        load_dotenv()
        api_key = os.getenv("GROQ_API_KEY")
        if not api_key:
             return "ERROR: ❌ GROQ_API_KEY missing."
             
        client = Groq(api_key=api_key)
        prompt = f"Summarize the following study text into 3-6 concise academic sentences:\n\n{text[:5000]}"

        import concurrent.futures
        from typing import Callable
        def _call() -> Any:
            response = client.chat.completions.create(
                model="llama-3.1-8b-instant",
                messages=[{"role": "user", "content": prompt}]
            )
            return response.choices[0].message.content if response.choices else None

        _call_fn: Callable[..., Any] = _call
        with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
            future = executor.submit(_call_fn)  # type: ignore[arg-type]
            return future.result(timeout=15)
    except Exception as e:
        return f"ERROR: ❌ {str(e)}"

def groq_safety_analysis(text):
    # Performs AI-driven safety scan to detect plagiarism and inappropriate content
    # Returns a structured JSON safety report
    try:
        api_key = os.getenv("GROQ_API_KEY")
        if not api_key:
            return {"overall": "Safe", "violation_type": "None", "reason": "AI safety skipped", "scores": {}}

        client = Groq(api_key=api_key)
        prompt = f"""Analyze this academic document text for integrity and safety. 
        You MUST return a valid JSON object with the following structure:
        {{
            "overall": "Safe" or "Flagged",
            "violation_type": "None" or the specific category,
            "reason": "Clear explanation of finding",
            "scores": {{
                "plagiarism": 0-10,
                "academic_misconduct": 0-10,
                "copyright_violation": 0-10,
                "spam_low_quality": 0-10,
                "inappropriate_content": 0-10,
                "malicious_content": 0-10
            }}
        }}

        Analyze these categories specifically:
        1. Plagiarism: Copied or suspicious text patterns.
        2. Academic Misconduct: Cheating guides, assignment selling, or unethical help.
        3. Copyright Violation: Sharing protected material without authorization.
        4. Spam / Low Quality: Promotional, garbled, or useless content.
        5. Inappropriate: Offensive, toxic, or adult material.
        6. Malicious: Phishing, malware instructions, or system harm.

        Text to analyze (limit 4000 chars):
        {text[:4000]}
        """

        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role": "system", "content": "You are an Academic Integrity Bot. Analyze text objectively and return only valid JSON."},
                      {"role": "user", "content": prompt}],
            temperature=0.1, # Low temperature for consistent JSON
            response_format={"type": "json_object"}
        )
        raw = response.choices[0].message.content.strip()
        return json.loads(raw)
    except Exception as e:
        return {"overall": "Safe", "violation_type": "None", "reason": f"AI error: {e}", "scores": {}}

def get_demo_response(tool_type: str, input_text: Optional[str]) -> str:
    # Generates human-readable fallback content when the Groq API is unavailable
    safe_input: str = str(input_text or "No content provided")
    preview: str = safe_input[0:20]  # type: ignore[index]
    responses = {
        'summary': "<p><strong>📝 Summary (Demo Mode):</strong> This is a fallback summary as the AI API is currently disconnected.</p>",
        'flashcards': '[{"q": "Demo Question?", "a": "Demo Answer."}]',
        'questions': "<p><strong>❓ Answer (Demo Mode):</strong> AI is in demo mode. Set your API key for real answers.</p>",
        'guide': f"<h4>Topic Study Guide (Demo)</h4><p>This is a structured guide template for {preview}...</p>"
    }
    return responses.get(tool_type, f"<p>Demo response for {tool_type}</p>")

# =====================================================
# SECTION: API Routes (Student Features)
# Purpose: HTTP endpoints for the student dashboard and tools
# =====================================================

@router.get("/dashboard", response_class=HTMLResponse)
async def dashboard(request: Request):
    # Displays the student's personal overview and recent activity
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    user_id = user["id"]
    latest_announcement = announcements_collection.find_one(
        {"$or": [{"audience": "all"}, {"audience": "student"}, {"audience": {"$exists": False}}]},
        sort=[("created_at", -1)]
    )
    
    return templates.TemplateResponse("student/dashboard.html", {
        "request": request,
        "user": user,
        "latest_announcement": latest_announcement,
        "total_uploads": notes_collection.count_documents({"uploader_id": user_id}),
        "approved_notes": notes_collection.count_documents({"uploader_id": user_id, "status": "approved"}),
        "pending_notes": notes_collection.count_documents({"uploader_id": user_id, "status": "pending"}),
        "saved_notes": favorites_collection.count_documents({"user_id": user_id}),
        "recent_activity": list(notifications_collection.find({"user_id": user_id}).sort("created_at", -1).limit(5)),
        **get_notif_context(user_id)
    })

@router.get("/upload", response_class=HTMLResponse)
async def upload_page(request: Request, retry: Optional[str] = None):
    # Renders the multi-file drag-and-drop upload interface
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    retry_note = None
    if retry:
        try:
            retry_note = notes_collection.find_one({"_id": ObjectId(retry), "uploader_id": user["id"]})
        except:
            pass

    db_subjects = list(subjects_collection.find({}, {"_id": 0, "name": 1}).sort("name", 1))
    subject_names = [s["name"] for s in db_subjects] or ["General"]
    
    return templates.TemplateResponse("student/upload.html", {
        "request": request,
        "user": user,
        "subjects": subject_names,
        "retry_note": retry_note,
        **get_notif_context(user["id"])
    })

@router.post("/upload")
async def upload_notes(request: Request, files: List[UploadFile] = File(...), titles: List[str] = Form(...), subjects: List[str] = Form(...)):
    # Handles note uploads, triggers AI safety scans, and notifies faculty
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    # Fetch global allowed file types from system settings
    # system_settings_collection is already imported from app.database
    settings = system_settings_collection.find_one({}) or {}
    allowed_exts = settings.get("allowed_file_types", "pdf,doc,docx,jpg,jpeg,png").split(',')
    allowed_exts = [ext.strip().lower() for ext in allowed_exts if ext.strip()]
    
    for i, file in enumerate(files):
        if not file.filename: continue
        
        # Security: Check file extension
        file_ext = file.filename.split('.')[-1].lower() if '.' in file.filename else ""
        if file_ext not in allowed_exts:
             continue # Skip unauthorized file types
        
        # Save file locally
        unique_filename = f"{user['id']}_{datetime.now(timezone.utc).timestamp()}_{file.filename}"
        file_path = os.path.join(UPLOAD_DIR, unique_filename)
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        # ── Step A: Extract text ─────────────────────────────────────────────
        extracted_text: str = ""
        fn = file.filename.lower()
        if fn.endswith('.pdf'):
            extracted_text = await asyncio.to_thread(extract_text_from_pdf, file_path)  # type: ignore[arg-type]
        elif any(fn.endswith(ext) for ext in ['.docx', '.pptx']):
            extracted_text = await asyncio.to_thread(extract_text_from_office, file_path)  # type: ignore[arg-type]

        note_title: str = titles[i] if i < len(titles) else file.filename
        note_subject: str = subjects[i] if i < len(subjects) else "General"

        # ── Step B: AI Moderation (harmful content) ──────────────────────────
        safety = await asyncio.to_thread(groq_safety_analysis, extracted_text) if extracted_text else {"overall": "Safe"}  # type: ignore[arg-type]
        ai_overall: str = safety.get("overall", "Safe")
        ai_violation: str = safety.get("violation_type", "None")
        ai_reason: str = safety.get("reason", "")
        ai_scores: dict = safety.get("scores", {})
        print(f"[AI Moderation] '{note_title}': {ai_overall} ({ai_violation})")

        # ── Step C: Duplicate & Plagiarism Detection ─────────────────────────
        plag_result = await asyncio.to_thread(
            check_duplicate_and_plagiarism,
            note_title, extracted_text, user["id"], note_subject
        )  # type: ignore[arg-type]
        plag_status: str = plag_result["status"]
        plag_score: float = plag_result["plagiarism_score"]
        plag_detail: str = plag_result["detail"]
        print(f"[Plagiarism] '{note_title}': {plag_status} (score={plag_score:.2%})")

        # ── Step D: Merge status — plagiarism overrides Safe, but NOT Flagged ─
        # Priority: Flagged > Possible Plagiarism > Duplicate Title > Safe
        if ai_overall == "Flagged":
            final_ai_status = "Flagged"
        elif plag_status == "Possible Plagiarism":
            final_ai_status = "Possible Plagiarism"
            ai_violation = "Plagiarism"
            ai_reason = plag_detail
        elif plag_status == "Duplicate Title":
            final_ai_status = "Duplicate Title"
            ai_violation = "Duplicate"
            ai_reason = plag_detail
        else:
            final_ai_status = ai_overall  # "Safe" or original

        # ── Step E: Save to database ─────────────────────────────────────────
        note_data = {
            "title": note_title,
            "subject": note_subject,
            "uploader_id": user["id"],
            "uploader_name": user["name"],
            "uploader_role": "student",
            "filename": unique_filename,
            "original_filename": file.filename,
            "file_path": file_path,
            "status": "pending",
            "uploaded_at": datetime.now(timezone.utc),
            "ai_status": final_ai_status,
            "ai_violation_type": ai_violation,
            "ai_flag_reason": ai_reason,
            "ai_detailed_results": ai_scores,
            "plagiarism_score": plag_score,
            "plagiarism_duplicate_of": plag_result.get("duplicate_of"),
            "plagiarism_duplicate_uploader": plag_result.get("duplicate_uploader"),
        }
        inserted = notes_collection.insert_one(note_data)

        # ── Step F: Auto-report to admin if plagiarism detected ──────────────
        if final_ai_status in ("Possible Plagiarism", "Duplicate Title", "Flagged"):
            try:
                reports_collection.insert_one({
                    "reported_item_id": str(inserted.inserted_id),
                    "item_type": "note",
                    "item_title": note_title,
                    "reporter_id": "system",
                    "reporter_name": "AI Detection System",
                    "reason": "PLAGIARISM" if final_ai_status in ("Possible Plagiarism", "Duplicate Title") else "INAPPROPRIATE_CONTENT",
                    "description": ai_reason or f"Automated detection: {final_ai_status}",
                    "status": "Open",
                    "ai_detected": True,
                    "created_at": datetime.now(timezone.utc)
                })
                # Notify all admins
                admin_users = list(users_collection.find({"role": "admin", "is_deleted": {"$ne": True}}))
                for adm in admin_users:
                    create_notification(
                        str(adm["_id"]), "SYSTEM",
                        "⚠️ AI Detection System",
                        f"AI detected issue in '{note_title}'"
                    )
                print(f"[Plagiarism] Auto-report created and admins notified.")
            except Exception as rpt_err:
                print(f"[Plagiarism] Auto-report failed: {rpt_err}")

        # ── Step G: Notify the student ───────────────────────────────────────
        if final_ai_status in ("Possible Plagiarism", "Duplicate Title"):
            create_notification(user["id"], "WARNING",
                "⚠️ Content Check Failed",
                f"Your note '{note_title}' was flagged: {plag_detail}")
        elif final_ai_status == "Flagged":
            create_notification(user["id"], "WARNING",
                "⚠️ Content Violation",
                f"Your note '{note_title}' was flagged for: {ai_violation}")
        else:
            create_notification(user["id"], "UPLOAD", "File Uploaded",
                f"'{note_title}' is pending review.")

    return RedirectResponse(url="/student/my-notes", status_code=303)

@router.get("/my-notes", response_class=HTMLResponse)
async def my_notes(request: Request):
    # Lists all notes uploaded by the current student with their status
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    notes = list(notes_collection.find({"uploader_id": user["id"]}).sort("uploaded_at", -1))
    return templates.TemplateResponse("student/my_notes.html", {
        "request": request, "user": user, "notes": notes, **get_notif_context(user["id"])
    })

@router.get("/browse", response_class=HTMLResponse)
async def browse_notes(request: Request):
    # Allows students to search and filter through approved community notes
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    notes = list(notes_collection.find({"status": "approved"}).sort("uploaded_at", -1))
    favs = [str(f["note_id"]) for f in favorites_collection.find({"user_id": user["id"]})]
    
    return templates.TemplateResponse("student/browse.html", {
        "request": request, "user": user, "notes": notes, "favorite_note_ids": favs, **get_notif_context(user["id"])
    })

@router.post("/favorite/{note_id}")
async def toggle_favorite(request: Request, note_id: str):
    # Bookmarks or removes a note from the student's saved collection
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    existing = favorites_collection.find_one({"user_id": user["id"], "note_id": note_id})
    if existing:
        favorites_collection.delete_one({"_id": existing["_id"]})
        return JSONResponse({"favorited": False})
    
    favorites_collection.insert_one({"user_id": user["id"], "note_id": note_id, "created_at": datetime.now(timezone.utc)})
    return JSONResponse({"favorited": True})

@router.get("/saved", response_class=HTMLResponse)
async def saved_notes(request: Request):
    # Renders the student's personal library of favorited notes
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    fav_ids = [ObjectId(f["note_id"]) for f in favorites_collection.find({"user_id": user["id"]})]
    notes = list(notes_collection.find({"_id": {"$in": fav_ids}, "status": "approved"}))
    return templates.TemplateResponse("student/saved.html", {
        "request": request, "user": user, "notes": notes, **get_notif_context(user["id"])
    })

@router.get("/preview/{note_id}", response_class=HTMLResponse)
async def preview_note(request: Request, note_id: str):
    # Visualizes note content (PDF/Image/Text) before downloading
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    note = notes_collection.find_one({"_id": ObjectId(note_id)})
    if not note: return HTMLResponse("Note not found", status_code=404)
    
    file_ext = os.path.splitext(note.get('filename', ''))[1].lower() if note.get('filename') else ""
    if not file_ext and note.get('original_filename'):
        file_ext = os.path.splitext(note.get('original_filename'))[1].lower()

    return templates.TemplateResponse("student/preview.html", {
        "request": request, 
        "user": user, 
        "note": note, 
        "file_ext": file_ext,
        **get_notif_context(user["id"])
    })

@router.post("/generate-summary/{note_id}")
async def fetch_summary(request: Request, note_id: str):
    # Generates an AI summary for the requested note
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    note = notes_collection.find_one({"_id": ObjectId(note_id)})
    if not note: return JSONResponse({"error": "Note not found"}, status_code=404)
    
    file_path = note.get("file_path")
    if not file_path or not os.path.exists(file_path):
        return JSONResponse({"error": "Original file not found for processing"}, status_code=404)
    
    # Extract text based on file type
    extracted_text: str = ""
    filename: str = str(note.get("filename", "")).lower()

    if filename.endswith(".pdf"):
        extracted_text = await asyncio.to_thread(extract_text_from_pdf, file_path)  # type: ignore[arg-type]
    elif any(filename.endswith(ext) for ext in [".docx", ".pptx"]):
        extracted_text = await asyncio.to_thread(extract_text_from_office, file_path)  # type: ignore[arg-type]
    elif any(filename.endswith(ext) for ext in [".jpg", ".jpeg", ".png"]):
        extracted_text = await asyncio.to_thread(extract_text_from_image, file_path)  # type: ignore[arg-type]
    elif any(filename.endswith(ext) for ext in [".txt", ".md"]):
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                extracted_text = f.read(5000)
        except Exception:
            pass

    if not extracted_text:
        return JSONResponse({"error": "Could not extract text from this file for summary."}, status_code=400)

    summary = await asyncio.to_thread(generate_ai_summary, extracted_text)  # type: ignore[arg-type]
    
    if summary and not summary.startswith("ERROR:"):
        return JSONResponse({"summary": summary})
    else:
        return JSONResponse({"error": summary or "AI service failed to generate summary."}, status_code=500)

@router.get("/download/{note_id}")
async def download_note(note_id: str):
    # Serves the physical file for local storage on the student's device
    note = notes_collection.find_one({"_id": ObjectId(note_id), "status": "approved"})
    if not note or not os.path.exists(note["file_path"]):
        return JSONResponse({"error": "File not found"}, status_code=404)
    
    return FileResponse(note["file_path"], filename=note.get("original_filename", "note.pdf"))

@router.get("/messages", response_class=HTMLResponse)
async def messages(request: Request):
    # Displays the student's inbox and active message threads
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    # Fetch all messages involving the student
    all_messages = list(messages_collection.find({
        "$or": [{"sender_id": str(user["id"])}, {"receiver_id": str(user["id"])}]
    }).sort("created_at", -1))
    
    conversations = {}
    for msg in all_messages:
        # Determine the other user in the conversation
        other_id = msg["receiver_id"] if msg["sender_id"] == str(user["id"]) else msg["sender_id"]
        
        # If we haven't added this conversation yet (it's the most recent message)
        if other_id not in conversations:
            other_user = users_collection.find_one({"_id": safe_id(other_id)})
            if other_user:
                conversations[other_id] = {
                    "user": other_user,
                    "last_message": msg
                }
    
    return templates.TemplateResponse("student/messages.html", {
        "request": request, 
        "user": user, 
        "conversations": conversations, 
        **get_notif_context(user["id"])
    })

@router.get("/chat/{other_user_id}", response_class=HTMLResponse)
async def chat_page(request: Request, other_user_id: str):
    # Renders a real-time-like chat interface with another user
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    other = users_collection.find_one({"_id": safe_id(other_user_id)})
    if not other: return HTMLResponse("User not found", status_code=404)
    
    # Fetch conversation history
    history = list(messages_collection.find({
        "$or": [
            {"sender_id": str(user["id"]), "receiver_id": str(other_user_id)}, 
            {"sender_id": str(other_user_id), "receiver_id": str(user["id"])}
        ]
    }).sort("created_at", 1))
    
    return templates.TemplateResponse("student/chat.html", {
        "request": request, 
        "user": user, 
        "other_user": other, 
        "chat_history": history, 
        **get_notif_context(user["id"])
    })

@router.post("/messages/send")
async def send_message(request: Request, receiver_id: str = Form(...), message: str = Form(...), subject: str = Form(None)):
    # Persists a new private message and alerts the receiver
    user = get_current_student(request)
    if not user: return JSONResponse({"success": False, "error": "Unauthorized"}, status_code=401)
    
    target = users_collection.find_one({"_id": safe_id(receiver_id)})
    if not target: return JSONResponse({"success": False, "error": "Receiver not found"}, status_code=404)
    
    msg_data = {
        "sender_id": str(user["id"]), 
        "sender_name": user.get("name", "Student"),
        "receiver_id": str(receiver_id), 
        "receiver_name": target.get("name", "User"),
        "subject": subject or "Direct Message",
        "message": message, 
        "is_read": False, 
        "created_at": datetime.now(timezone.utc)
    }
    
    messages_collection.insert_one(msg_data)
    create_notification(str(receiver_id), "MESSAGE", "New Message", f"From {user['name']}")
    return JSONResponse({"success": True})

@router.get("/forum", response_class=HTMLResponse)
async def forum(request: Request):
    # Academic discussion board where students can post questions
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    posts = list(forum_posts_collection.find().sort("created_at", -1))
    return templates.TemplateResponse("student/forum.html", {
        "request": request, "user": user, "posts": posts, **get_notif_context(str(user["id"]))
    })

@router.post("/forum/post")
async def create_forum_post(request: Request, title: str = Form(...), subject: str = Form(...), content: str = Form(...)):
    # Creates a new forum thread for community discussion
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    post_data = {"title": title, "subject": subject, "content": content, "author_id": str(user["id"]), "author_name": user["name"], "replies": [], "created_at": datetime.now(timezone.utc)}
    forum_posts_collection.insert_one(post_data)
    return RedirectResponse(url="/student/forum", status_code=303)

@router.post("/forum/reply/{post_id}")
async def reply_to_forum_post(request: Request, post_id: str, reply: str = Form(...)):
    # Connects a student's reply to an existing forum thread
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    from bson import ObjectId
    try: oid = ObjectId(post_id)
    except: return JSONResponse({"error": "Invalid post ID"}, status_code=400)
    
    reply_data = {
        "author_id": str(user["id"]),
        "author_name": user["name"],
        "author_role": "student",
        "content": reply,
        "created_at": datetime.now(timezone.utc)
    }
    
    forum_posts_collection.update_one(
        {"_id": oid},
        {"$push": {"replies": reply_data}}
    )
    
    # Notify original author (if it's not the same user)
    post = forum_posts_collection.find_one({"_id": oid})
    if post and str(post.get("author_id")) != str(user["id"]):
        create_notification(
            post["author_id"], "FORUM_REPLY", 
            "New Forum Reply", 
            f"{user['name']} replied to your discussion: '{post.get('title')}'"
        )
        
    return JSONResponse({"success": True})

@router.get("/ai", response_class=HTMLResponse)
async def ai_assistant(request: Request):
    # Landing page for AI study tools (Summary, Flashcards, Q&A)
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    return templates.TemplateResponse("student/ai_assistant.html", {
        "request": request, "user": user, **get_notif_context(str(user["id"]))
    })

@router.post("/ai/process")
async def process_ai_request(request: Request, tool_type: str = Form(...), input_text: str = Form(...)):
    # Routes AI requests to Groq (Llama) and returns educational aids
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    try:
        api_key = os.getenv("GROQ_API_KEY")
        if not api_key: return JSONResponse({"result": get_demo_response(tool_type, input_text), "ai_powered": False})
        
        client = Groq(api_key=api_key)
        
        # Define specific prompts and models for each tool (aligned with dashboard output sizes)
        prompts = {
            "summary": "You are an Academic Assistant. Generate a SHORT summary using 3-5 concise, high-impact bullet points maximum.",
            "flashcards": "You are a Study Assistant. Create exactly 5 flashcards for the topic provided. Return ONLY a valid JSON list of objects. Each object MUST have 'q' (question) and 'a' (answer) keys. Each answer must be a single complete sentence. Do not repeat concepts. Format example: [{\"q\": \"...\", \"a\": \"...\"}, ...]",
            "questions": "You are an Expert Tutor. Provide an accurate and helpful answer to the question. The length should VARY based on the complexity of the query.",
            "guide": "You are a Study Guide Architect. Create a MEDIUM-TO-LONG, comprehensive study guide. Include detailed sections, definitions, and key concepts."
        }
        
        system_prompt = prompts.get(tool_type, "You are a helpful Academic Assistant.")
        # Use 8B Instant for all tasks as per performance optimization
        model = "llama-3.1-8b-instant"
        
        # Use JSON mode for flashcards to ensure front-end parsing works
        response_kwargs = {}
        if tool_type == "flashcards":
            response_kwargs["response_format"] = {"type": "json_object"}
            
        resp = client.chat.completions.create(
            model=model, 
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": input_text}
            ],
            **response_kwargs
        )
        
        result = resp.choices[0].message.content
        
        # If flashcards, ensure it's a list (JSON object mode might wrap it in an object)
        if tool_type == "flashcards" and result:
            try:
                raw_result: str = result if isinstance(result, str) else ""
                data = json.loads(raw_result)
                # If AI returned {"flashcards": [...]}, extract the list
                if isinstance(data, dict):
                    original_result = raw_result
                    # Check common keys safely
                    for key in ["flashcards", "cards", "data"]:
                        val = data.get(key)
                        if val and isinstance(val, list):
                            result = json.dumps(val)
                            break

                    # Fallback: if data itself is a dict with mixed content
                    if result == original_result:
                        for v in data.values():
                            if isinstance(v, list):
                                result = json.dumps(v)
                                break
            except Exception:
                pass

        return JSONResponse({"success": True, "result": result, "ai_powered": True})
    except Exception as e:
        print(f"AI Process Error: {e}")
        return JSONResponse({"result": get_demo_response(tool_type, input_text), "ai_powered": False, "error": str(e)})

@router.get("/announcements", response_class=HTMLResponse)
async def announcements_page(request: Request):
    # Displays all campus-wide and student-specific announcements
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    all_announcements = list(announcements_collection.find({
        "$or": [{"audience": "all"}, {"audience": "student"}, {"audience": {"$exists": False}}]
    }).sort("created_at", -1))
    
    return templates.TemplateResponse("student/announcements.html", {
        "request": request, "user": user, "announcements": all_announcements, **get_notif_context(user["id"])
    })

@router.get("/notifications", response_class=HTMLResponse)
async def notifications_page(request: Request):
    # Renders the full notification inbox for the student
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    notifications = list(notifications_collection.find({"user_id": user["id"]}).sort("created_at", -1))
    return templates.TemplateResponse("student/notifications.html", {
        "request": request, "user": user, "notifications": notifications, **get_notif_context(user["id"])
    })

@router.post("/notifications/mark-read")
async def mark_notification_read(request: Request, notif_id: str = Form(...)):
    # Marks a specific notification as viewed by the student (form-field version)
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    try:
        notifications_collection.update_one(
            {"_id": ObjectId(notif_id), "user_id": user["id"]},
            {"$set": {"is_read": True}}
        )
    except Exception:
        pass
    return JSONResponse({"status": "Success"})

@router.post("/notifications/mark-read/{notif_id}")
async def mark_notification_read_path(request: Request, notif_id: str):
    # Marks a specific notification as read (path-param version used by frontend JS)
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    try:
        notifications_collection.update_one(
            {"_id": ObjectId(notif_id), "user_id": user["id"]},
            {"$set": {"is_read": True}}
        )
    except Exception:
        pass
    return JSONResponse({"status": "Success"})

@router.post("/notifications/mark-all-read")
async def mark_all_notifications_read(request: Request):
    # Marks all unread alerts as read for the current student
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    notifications_collection.update_many(
        {"user_id": user["id"], "is_read": False},
        {"$set": {"is_read": True}}
    )
    return JSONResponse({"status": "Success"})

@router.delete("/notifications/clear-all")
async def clear_all_notifications(request: Request):
    # Permanently deletes all notifications for the current student
    user = get_current_student(request)
    if not user: return JSONResponse({"error": "Unauthorized"}, status_code=401)
    
    result = notifications_collection.delete_many({"user_id": user["id"]})
    return JSONResponse({"status": "Success", "deleted": result.deleted_count})

@router.get("/faculty", response_class=HTMLResponse)
async def faculty_directory(request: Request):
    # Lists all active faculty members for student contact and messaging
    user = get_current_student(request)
    if not user: return RedirectResponse(url="/login", status_code=303)
    
    faculty_list = list(users_collection.find({"role": "faculty", "is_deleted": {"$ne": True}}))
    db_subjects = list(subjects_collection.find().sort("name", 1))
    
    return templates.TemplateResponse("student/faculty.html", {
        "request": request, 
        "user": user, 
        "faculty_list": faculty_list, 
        "subjects": db_subjects,
        **get_notif_context(user["id"])
    })
